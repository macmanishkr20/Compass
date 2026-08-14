"""Exporting a design.

A design is one standalone HTML document, so every format here is produced by
rendering that document in headless Chromium — the same engine the preview
iframe uses, which is what makes the export match what the user sees.

    html   the document itself (served straight from the store)
    pdf    Chromium's print output
    png    a full-page screenshot
    zip    the document plus its design-system notes and a README
    pptx   one PowerPoint slide per marked slide, each a full-bleed render

PPTX needs python-pptx; PDF and PNG need Playwright. Both are optional — a
missing one raises RuntimeError, which the API turns into a 501 rather than
pretending the format exists.
"""

from __future__ import annotations

import io
import zipfile

_NO_PLAYWRIGHT = (
    "Playwright is not installed on the server host "
    "(pip install playwright && playwright install chromium)."
)
_NO_PPTX = "python-pptx is not installed on the server host (pip install python-pptx)."

# What marks one slide. Deliberately explicit — a bare `section` would turn an
# ordinary long page into a stack of half-height pages, so the slides template
# tells the model to mark each slide with class="slide" instead of guessing.
_SLIDE_SELECTORS = ("[data-slide]", ".slide")


# A deck that presents itself shows one slide at a time — the others are
# hidden, stacked or translated off-screen by its own script. Exporting has to
# see all of them, so the whole deck is unfolded first and any chrome the
# author marked as presentation-only is dropped.
_UNFOLD_SLIDES = """
  .slide, [data-slide] {
    display: block !important;
    visibility: visible !important;
    opacity: 1 !important;
    position: static !important;
    transform: none !important;
    inset: auto !important;
  }
  [class*="deck"], [class*="slides"], [id*="deck"], [id*="slides"] {
    height: auto !important;
    max-height: none !important;
    overflow: visible !important;
  }
  [data-export-hide], .deck-nav, .slide-nav, .deck-controls { display: none !important; }
"""


async def _render(html: str, width: int, height: int, *, scale: float = 1):
    """Open the document in Chromium. Caller drives the page, then closes it."""
    try:
        from playwright.async_api import async_playwright
    except ImportError as err:  # pragma: no cover - host without Playwright
        raise RuntimeError(_NO_PLAYWRIGHT) from err

    pw = await async_playwright().start()
    browser = await pw.chromium.launch(args=["--no-sandbox"])
    page = await browser.new_page(
        viewport={"width": width, "height": height}, device_scale_factor=scale
    )
    # set_content rather than a data: URL — @import'd fonts need a real origin
    # to resolve against, and about:blank gives them one.
    await page.set_content(html, wait_until="load")
    await page.wait_for_timeout(700)  # let webfonts and entry animations settle

    async def close() -> None:
        await browser.close()
        await pw.stop()

    return page, close


async def _unfold(page) -> int:
    """Show every slide at once. Returns how many the document has."""
    count = 0
    for selector in _SLIDE_SELECTORS:
        found = await page.query_selector_all(selector)
        if len(found) > 1:
            count = len(found)
            break
    if count:
        await page.add_style_tag(content=_UNFOLD_SLIDES)
        await page.wait_for_timeout(250)  # let the relayout settle
    return count


async def _slide_box(page) -> dict | None:
    """The page box a deck should print at, or None if this isn't a deck.

    A deck built from fixed 16:9 sections prints one slide per page at that
    size. A deck whose slides are whatever height their content needs — which
    is most of them once unfolded — has to print at the tallest, or the long
    ones split across pages and the short ones leave a gap.
    """
    for selector in _SLIDE_SELECTORS:
        slides = await page.query_selector_all(selector)
        if len(slides) <= 1:
            continue
        boxes = [b for b in [await s.bounding_box() for s in slides] if b]
        if not boxes:
            continue
        widest = max(b["width"] for b in boxes)
        tallest = max(b["height"] for b in boxes)
        return {"width": widest, "height": tallest + 2}
    return None


async def to_pdf(html: str, *, width: int = 1280) -> bytes:
    page, close = await _render(html, width, 900)
    try:
        # print_background keeps the palette; the design decides its own size,
        # so the page box follows the rendered width rather than a paper size.
        await _unfold(page)
        box = await _slide_box(page)
        if box:
            # A deck prints one slide per page. Without the break rule Chromium
            # would flow the slides together and cut them mid-height.
            await page.add_style_tag(
                content=(
                    "@media print{"
                    + ",".join(_SLIDE_SELECTORS)
                    + "{break-after:page;break-inside:avoid}}"
                )
            )
            page_width, page_height = box["width"], box["height"]
        else:
            page_width = width
            # One tall page, so nothing is cut mid-section. scrollHeight alone
            # under-measures when the body's own box is the taller one, which
            # spills the footer onto a second page.
            page_height = await page.evaluate(
                "Math.ceil(Math.max("
                "document.documentElement.scrollHeight,"
                "document.body.scrollHeight,"
                "document.body.getBoundingClientRect().bottom)) + 2"
            )
        return await page.pdf(
            print_background=True,
            width=f"{page_width}px",
            height=f"{page_height}px",
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
        )
    finally:
        await close()


async def to_png(html: str, *, width: int = 1280) -> bytes:
    # Twice the device scale: the export is a picture of the design, and a
    # 1x screenshot of a page is soft the moment anyone zooms it.
    page, close = await _render(html, width, 900, scale=2)
    try:
        await _unfold(page)
        return await page.screenshot(full_page=True, type="png")
    finally:
        await close()


async def to_thumbnail(html: str) -> bytes:
    """A small PNG of the top of the design, for the projects table.

    Rendered at full width and scaled down by the device pixel ratio rather
    than resized afterwards — that keeps it sharp without needing an imaging
    library on the host.
    """
    try:
        from playwright.async_api import async_playwright
    except ImportError as err:  # pragma: no cover - host without Playwright
        raise RuntimeError(_NO_PLAYWRIGHT) from err

    pw = await async_playwright().start()
    browser = await pw.chromium.launch(args=["--no-sandbox"])
    try:
        page = await browser.new_page(
            viewport={"width": 1280, "height": 1600}, device_scale_factor=0.25
        )
        await page.set_content(html, wait_until="load")
        await page.wait_for_timeout(500)
        return await page.screenshot(type="png")
    finally:
        await browser.close()
        await pw.stop()


def to_zip(*, name: str, html: str, prompt: str, system_notes: str = "") -> bytes:
    readme = [
        f"# {name}",
        "",
        "Designed with Compass Design.",
        "",
        "## Brief",
        "",
        prompt or "(no brief recorded)",
        "",
        "## Files",
        "",
        "- `index.html` — the design. Open it in any browser; it is self-contained.",
    ]
    if system_notes:
        readme.append("- `design-system.md` — the system this design follows.")

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("index.html", html)
        z.writestr("README.md", "\n".join(readme) + "\n")
        if system_notes:
            z.writestr("design-system.md", system_notes)
    return buf.getvalue()


async def to_pptx(html: str, *, width: int = 1600) -> bytes:
    """Render each slide to an image and lay them out as a deck.

    HTML and PowerPoint share no layout model, so a faithful export means
    rendering rather than translating. Two things decide whether the result
    looks like the design or like a photocopy of it:

    * the deck's page is sized to the slides' own aspect, so nothing is
      stretched to fit a 16:9 box it was never drawn in; and
    * the render happens at twice the device scale, because a slide laid out
      at 700-odd CSS pixels stretched across thirteen inches is about 50 DPI
      — which is exactly what "pixelated" looks like.

    A slide that doesn't share the deck's aspect is fitted inside it and
    centred on the design's own background rather than distorted.
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Emu
    except ImportError as err:  # pragma: no cover - host without python-pptx
        raise RuntimeError(_NO_PPTX) from err

    try:
        from playwright.async_api import async_playwright
    except ImportError as err:  # pragma: no cover - host without Playwright
        raise RuntimeError(_NO_PLAYWRIGHT) from err

    pw = await async_playwright().start()
    browser = await pw.chromium.launch(args=["--no-sandbox"])
    try:
        page = await browser.new_page(
            viewport={"width": width, "height": 900}, device_scale_factor=2
        )
        await page.set_content(html, wait_until="load")
        await page.wait_for_timeout(700)
        await _unfold(page)

        background = await page.evaluate(
            "getComputedStyle(document.body).backgroundColor || 'rgb(255,255,255)'"
        )

        shots: list[tuple[bytes, float]] = []   # (png, aspect)
        for selector in _SLIDE_SELECTORS:
            slides = await page.query_selector_all(selector)
            if len(slides) > 1:
                for slide in slides:
                    box = await slide.bounding_box()
                    if not box or box["width"] < 2 or box["height"] < 2:
                        continue
                    try:
                        shots.append(
                            (await slide.screenshot(type="png"), box["width"] / box["height"])
                        )
                    except Exception:  # noqa: BLE001 - a slide with no box
                        continue
                break
        if not shots:  # not a deck — one slide holding the whole design
            box = await page.evaluate(
                "[document.documentElement.scrollWidth, document.documentElement.scrollHeight]"
            )
            shots = [
                (
                    await page.screenshot(full_page=True, type="png"),
                    (box[0] / box[1]) if box[1] else 16 / 9,
                )
            ]
    finally:
        await browser.close()
        await pw.stop()

    # A deck drawn to one shape keeps that shape. A deck whose slides are
    # whatever height their content needs — which is most of them once
    # unfolded — gets the standard 16:9 canvas, and each slide is fitted
    # inside it rather than forced into it.
    aspects = [a for _, a in shots]
    uniform = max(aspects) - min(aspects) <= min(aspects) * 0.05
    aspect = aspects[0] if uniform else 16 / 9

    deck = Presentation()
    deck.slide_width = Emu(12192000)                    # 13.333in, the usual canvas
    deck.slide_height = Emu(int(12192000 / aspect))
    blank = deck.slide_layouts[6]
    fill = _rgb(background)

    for png, aspect in shots:
        slide = deck.slides.add_slide(blank)
        if fill:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*fill)
        # Contain, never stretch: a slide of a different shape is centred.
        page_w, page_h = deck.slide_width, deck.slide_height
        w = page_w
        h = int(page_w / aspect)
        if h > page_h:
            h = page_h
            w = int(page_h * aspect)
        slide.shapes.add_picture(
            io.BytesIO(png), int((page_w - w) / 2), int((page_h - h) / 2), width=w, height=h
        )

    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


def _rgb(css: str) -> tuple[int, int, int] | None:
    """The r,g,b of a computed CSS colour, if it is opaque enough to matter."""
    import re

    m = re.match(r"rgba?\(([^)]+)\)", css or "")
    if not m:
        return None
    parts = [p.strip() for p in m.group(1).replace("/", ",").split(",")]
    try:
        r, g, b = (int(float(parts[i])) for i in range(3))
    except (ValueError, IndexError):
        return None
    if len(parts) > 3:
        try:
            if float(parts[3]) < 0.5:
                return None
        except ValueError:
            pass
    return r, g, b
